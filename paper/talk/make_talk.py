#!/usr/bin/env python3
"""Build the EUROMECH talk deck (editable .pptx) for the vortex-gust JEPA paper.

Figures are pre-rasterized PNGs in paper/talk/figs/ (see the conversion step in
the build notes). Numbers come from paper/HEADLINE_NUMBERS.md.

Run:  python paper/talk/make_talk.py   ->   paper/talk/euromech_gust_jepa.pptx
"""
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

HERE = Path(__file__).resolve().parent
FIG = HERE / "figs"
OUT = HERE / "euromech_gust_jepa.pptx"

# ---- palette / type ---------------------------------------------------------
NAVY = RGBColor(0x1F, 0x38, 0x64)
ACCENT = RGBColor(0xC0, 0x52, 0x10)   # warm rust for emphasis / glyphs
INK = RGBColor(0x24, 0x29, 0x33)
MUTE = RGBColor(0x6B, 0x72, 0x80)
RULE = RGBColor(0xD6, 0xDB, 0xE3)
LIGHT = RGBColor(0xC9, 0xD3, 0xE6)
FONT = "Calibri"

W, H = 13.333, 7.5

prs = Presentation()
prs.slide_width = Inches(W)
prs.slide_height = Inches(H)
BLANK = prs.slide_layouts[6]


# ---- low-level helpers ------------------------------------------------------
def slide():
    return prs.slides.add_slide(BLANK)


def rect(s, l, t, w, h, fill=None, line=None):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1.0)
    return sp


def _runs(p, text, size, color, base_bold=False, bold_color=NAVY):
    for i, seg in enumerate(text.split("**")):
        if seg == "":
            continue
        r = p.add_run()
        r.text = seg
        r.font.name = FONT
        r.font.size = Pt(size)
        emph = (i % 2 == 1)
        r.font.bold = base_bold or emph
        r.font.color.rgb = bold_color if emph else color


def textbox(s, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tf


def para(tf, text, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT,
         space_before=0, space_after=6, bold_color=None, first=False, glyph=None):
    if first and not tf.paragraphs[0].runs:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    if glyph:
        g = p.add_run()
        g.text = glyph
        g.font.name = FONT
        g.font.size = Pt(size)
        g.font.bold = True
        g.font.color.rgb = bold_color or ACCENT
    _runs(p, text, size, color, base_bold=bold, bold_color=bold_color or NAVY)
    return p


def bullets(s, items, l, t, w, h, base=17):
    tf = textbox(s, l, t, w, h)
    first = True
    for it in items:
        lvl = it.get("lvl", 0)
        sz = it.get("sz", base if lvl == 0 else base - 3)
        sa = it.get("sa", 9 if lvl == 0 else 5)
        if lvl == 0:
            para(tf, it["t"], size=sz, color=INK, space_after=sa, glyph="•  ", first=first)
        else:
            para(tf, "      –  " + it["t"], size=sz, color=MUTE, space_after=sa, first=first)
        first = False
    return tf


def fig(s, name, l, t, w, h, halign="center", valign="middle"):
    path = FIG / name
    iw, ih = Image.open(path).size
    box_ar, img_ar = w / h, iw / ih
    if img_ar > box_ar:
        nw, nh = w, w / img_ar
    else:
        nh, nw = h, h * img_ar
    x = l + (w - nw) / 2 if halign == "center" else (l if halign == "left" else l + (w - nw))
    y = t + (h - nh) / 2 if valign == "middle" else (t if valign == "top" else t + (h - nh))
    s.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(nw), Inches(nh))


ASSET = HERE / "assets"
_ASSET_CACHE = FIG / "_assetcache"
ASSETS = {
    "euromech": "IMG_2285.jpeg", "inta": "IMG_2286.png", "uc3m": "IMG_2289.png",
    "bsc": "IMG_2287.png", "upc": "IMG_2288.png", "bbva": "IMG_2283.jpeg",
    "redleonardo": "IMG_2282.png", "photo_solera": "IMG_2292.png",
    "photo_miro": "IMG_2290.jpeg", "photo_lehmkuhl": "IMG_2291.jpeg",
}


def asset_path(key):
    p = ASSET / ASSETS[key]
    return p if p.exists() else None


def place_image(s, path, l, t, w, h, halign="center", valign="middle"):
    iw, ih = Image.open(path).size
    box_ar, img_ar = w / h, iw / ih
    if img_ar > box_ar:
        nw, nh = w, w / img_ar
    else:
        nh, nw = h, h * img_ar
    x = l + (w - nw) / 2 if halign == "center" else (l if halign == "left" else l + (w - nw))
    y = t + (h - nh) / 2 if valign == "middle" else (t if valign == "top" else t + (h - nh))
    s.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(nw), Inches(nh))


def logo(s, key, l, t, w, h, halign="center", valign="middle", crop=None):
    """Aspect-fit a logo (by asset key) into a box; silently skip if missing."""
    p = asset_path(key)
    if p is None:
        return
    if crop is not None:
        _ASSET_CACHE.mkdir(exist_ok=True)
        p = _ASSET_CACHE / f"{key}_crop.png"
        Image.open(asset_path(key)).crop(crop).save(p)
    place_image(s, p, l, t, w, h, halign, valign)


def square_photo(key, focus_x=0.5, focus_y=0.45):
    """Center-crop a photo to a square (with focus bias), cache, return path or None."""
    p = asset_path(key)
    if p is None:
        return None
    _ASSET_CACHE.mkdir(exist_ok=True)
    out = _ASSET_CACHE / f"{key}_sq.png"
    im = Image.open(p).convert("RGB")
    w, h = im.size
    side = min(w, h)
    cx, cy = int(w * focus_x), int(h * focus_y)
    left = max(0, min(w - side, cx - side // 2))
    top = max(0, min(h - side, cy - side // 2))
    im.crop((left, top, left + side, top + side)).save(out)
    return str(out)


def caption(s, text, l, t, w):
    tf = textbox(s, l, t, w, 0.4)
    para(tf, text, size=11, color=MUTE, align=PP_ALIGN.CENTER, first=True)


def header(s, title, kicker=None):
    if kicker:
        tf = textbox(s, 0.55, 0.30, 11.5, 0.32)
        para(tf, kicker.upper(), size=12.5, color=ACCENT, bold=True, space_after=0, first=True)
    tf2 = textbox(s, 0.55, 0.60, 12.25, 0.95)
    para(tf2, title, size=26, color=NAVY, bold=True, space_after=0, first=True)
    rect(s, 0.55, 1.50, 12.23, 0.022, fill=RULE)
    rect(s, 0.55, 1.485, 2.1, 0.05, fill=ACCENT)


def left_footer(s):
    tf = textbox(s, 0.55, 7.06, 9.5, 0.3)
    para(tf, "Sparse-sensing state estimation, vortex-gust interactions  ·  EUROMECH colloquium",
         size=9.5, color=MUTE, first=True)


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


# ---- slide templates --------------------------------------------------------
def s_title():
    s = slide()
    rect(s, 0, 0, 0.30, H, fill=NAVY)
    tf = textbox(s, 0.95, 1.35, 11.7, 2.5)
    para(tf, "Sparse-sensing state estimation of vortex-gust airfoil interactions",
         size=31, color=NAVY, bold=True, space_after=6, first=True)
    para(tf, "Toward model-based flow control", size=21, color=ACCENT, bold=True, space_after=0)
    rect(s, 1.0, 3.95, 3.0, 0.06, fill=ACCENT)
    tf2 = textbox(s, 1.0, 4.2, 11.5, 1.5)
    para(tf2, "A. Solera-Rico, A. Miró, O. Lehmkuhl, **C. Sanmiguel Vila**",
         size=15, color=INK, first=True, space_after=3)
    para(tf2, "csanvil@inta.es",
         size=12.5, color=ACCENT, space_after=0)
    # conference logo (top-right) and institution logo strip (bottom)
    logo(s, "euromech", 11.55, 0.42, 1.15, 0.72, crop=(0, 0, 460, 282))
    _ly = 5.35
    logo(s, "inta", 1.0, _ly, 0.66, 0.66, halign="left")
    logo(s, "uc3m", 1.95, _ly + 0.08, 1.55, 0.5, halign="left")
    logo(s, "bsc", 3.75, _ly, 0.74, 0.66, halign="left")
    logo(s, "upc", 4.72, _ly, 0.66, 0.66, halign="left")
    tf3 = textbox(s, 1.0, 6.55, 11.5, 0.5)
    para(tf3, "EUROMECH Colloquium  ·  Data-driven active control in flows: from model-based to reinforcement learning",
         size=12, color=NAVY, bold=True, first=True)
    notes(s, "Frame: extreme gust control needs a fast, faithful forward model. Our work "
              "asks which reduced state is worth planning against, and shows a predictive "
              "(JEPA) latent captures the wake physics that reconstructive and linear "
              "baselines lose. Audience is new to JEPA, so I will spend a slide on what it is.")
    return s


def s_bullets(title, kicker, items, note="", lead=None, base=18):
    s = slide()
    header(s, title, kicker)
    if lead:
        tf = textbox(s, 0.6, 1.78, 12.1, 0.8)
        para(tf, lead, size=19, color=INK, first=True)
        bullets(s, items, 0.7, 2.65, 12.0, 4.1, base=base)
    else:
        bullets(s, items, 0.7, 1.9, 12.0, 5.0, base=base)
    left_footer(s)
    if note:
        notes(s, note)
    return s


def s_fig_right(title, kicker, items, figname, cap=None, note="", fl=7.7, fw=5.2, base=17):
    s = slide()
    header(s, title, kicker)
    bullets(s, items, 0.6, 1.85, fl - 0.85, 5.0, base=base)
    fig(s, figname, fl, 1.82, fw, 4.8, valign="top")
    if cap:
        caption(s, cap, fl, 6.68, fw)
    left_footer(s)
    if note:
        notes(s, note)
    return s


def s_fig_below(title, kicker, items, figname, cap=None, note="", ft=3.55, fh=3.05, base=17):
    s = slide()
    header(s, title, kicker)
    bullets(s, items, 0.7, 1.78, 12.0, ft - 1.85, base=base)
    fig(s, figname, 0.8, ft, 11.7, fh, valign="top")
    if cap:
        caption(s, cap, 0.8, ft + fh + 0.02, 11.7)
    left_footer(s)
    if note:
        notes(s, note)
    return s


def s_question(title, kicker, items, quote, note=""):
    s = slide()
    header(s, title, kicker)
    bullets(s, items, 0.7, 1.85, 12.0, 2.4, base=18)
    rect(s, 0.7, 4.55, 11.9, 1.7, fill=RGBColor(0xF3, 0xF5, 0xF9))
    rect(s, 0.7, 4.55, 0.10, 1.7, fill=ACCENT)
    tf = textbox(s, 1.05, 4.55, 11.3, 1.7, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, quote, size=22, color=NAVY, bold=True, first=True)
    left_footer(s)
    if note:
        notes(s, note)
    return s


def s_divider(title, sub=None):
    s = slide()
    rect(s, 0, 0, W, H, fill=NAVY)
    tf = textbox(s, 1.0, 2.85, 11.3, 1.6, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, title, size=36, color=RGBColor(0xFF, 0xFF, 0xFF), bold=True, first=True)
    if sub:
        para(tf, sub, size=16, color=LIGHT, space_before=6)
    rect(s, 1.05, 4.25, 3.0, 0.06, fill=ACCENT)
    return s


def b0(t, **k):
    return {"t": t, "lvl": 0, **k}


def b1(t, **k):
    return {"t": t, "lvl": 1, **k}


# ---- native schematic helpers (large, editable; no tiny raster text) --------
DGRAY = RGBColor(0x8A, 0x93, 0xA3)
AE_RED = RGBColor(0xB0, 0x3A, 0x2E)
AE_FILL = RGBColor(0xF7, 0xEC, 0xEA)
JP_GRN = RGBColor(0x2E, 0x7D, 0x4F)
JP_FILL = RGBColor(0xEA, 0xF2, 0xEC)
BOX_FILL = RGBColor(0xEE, 0xF1, 0xF7)
COND_FILL = RGBColor(0xFD, 0xF0, 0xD9)


def dbox(s, l, t, w, h, text, fill=BOX_FILL, edge=NAVY, tcolor=INK, size=14):
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.shadow.inherit = False
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = edge
    sp.line.width = Pt(1.5)
    tf = sp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(3)
    tf.margin_right = Pt(3)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    for j, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        _runs(p, line, size, tcolor, base_bold=True)
    return sp


def ashape(s, kind, l, t, w, h, fill=DGRAY):
    sp = s.shapes.add_shape(kind, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.shadow.inherit = False
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    sp.line.fill.background()
    return sp


def harrow(s, l, t, w, h=0.18, fill=DGRAY):
    return ashape(s, MSO_SHAPE.RIGHT_ARROW, l, t, w, h, fill)


def uparrow(s, l, t, w, h, fill=DGRAY):
    return ashape(s, MSO_SHAPE.UP_ARROW, l, t, w, h, fill)


def dlabel(s, l, t, w, text, size=12.5, color=MUTE, align=PP_ALIGN.LEFT, bold=False):
    tf = textbox(s, l, t, w, 0.5)
    para(tf, text, size=size, color=color, align=align, bold=bold, first=True)


def slide_jepa_concept():
    s = slide()
    header(s, "What is a JEPA?  (1 / 2)", "Method, primer")
    bx = [0.9, 3.15, 5.5, 7.45, 9.9]
    bw = [1.7, 1.7, 1.3, 1.8, 2.55]
    ax = [2.6, 4.85, 6.8, 9.35]
    hb = 0.72
    # Row A: autoencoder (reconstruct)
    dlabel(s, 0.9, 1.74, 11.8, "Autoencoder (most flow ROMs): reconstruct the field",
           size=15, color=AE_RED, bold=True)
    ya = 2.22
    ae = ["field ω(t)", "Encoder", "latent z", "Decoder", "reconstructed field"]
    for i, txt in enumerate(ae):
        dbox(s, bx[i], ya, bw[i], hb, txt, fill=AE_FILL, edge=AE_RED, size=13.5)
    for axx in ax:
        harrow(s, axx, ya + hb / 2 - 0.09, 0.5)
    dlabel(s, 0.9, 3.06, 11.8, "loss compares the reconstruction to the input in PIXEL space",
           size=13, color=AE_RED)
    # Row B: JEPA (predict)
    dlabel(s, 0.9, 4.04, 11.8, "JEPA: predict the next latent, no decoder",
           size=15, color=JP_GRN, bold=True)
    yb = 4.52
    jp = ["field ω(t)", "Encoder", "latent z(t)", "Predictor", "predicted z(t+1)"]
    for i, txt in enumerate(jp):
        dbox(s, bx[i], yb, bw[i], hb, txt, fill=JP_FILL, edge=JP_GRN, size=13.5)
    for axx in ax:
        harrow(s, axx, yb + hb / 2 - 0.09, 0.5)
    dlabel(s, 0.9, 5.36, 11.8,
           "loss compares predicted z(t+1) to z(t+1) = Encoder(next frame) in LATENT space; "
           "an anti-collapse term replaces the decoder", size=13, color=JP_GRN)
    tf = textbox(s, 0.9, 6.18, 11.8, 0.6)
    para(tf, "**Reconstruction ties the latent to pixels; prediction makes it forecastable**, "
             "which is what a controller needs.", size=15, color=NAVY, first=True)
    left_footer(s)
    notes(s, "Pedagogic core for a non-ML audience. Same first three boxes; the only change is "
             "Decoder -> Predictor and where the loss lives: pixels (AE) vs latent space (JEPA). "
             "No decoder in JEPA; the anti-collapse term stops the encoder collapsing.")
    return s


def slide_our_arch():
    s = slide()
    header(s, "Our architecture for gust encounters", "Method")
    ya = 1.95
    hb = 0.95
    bx = [0.6, 2.85, 5.15, 7.05, 9.35]
    bw = [2.0, 2.0, 1.6, 1.95, 2.0]
    arx = [2.6, 4.85, 6.65, 9.05]
    boxes = ["mid-plane\nvorticity ω(t)", "Encoder\n(CNN + ViT)",
             "latent z(t)\nd = 32 / 64", "Predictor\n(AR transformer)", "predicted\nz(t+1)"]
    fills = [BOX_FILL, BOX_FILL, JP_FILL, JP_FILL, JP_FILL]
    edges = [NAVY, NAVY, JP_GRN, JP_GRN, JP_GRN]
    for i, txt in enumerate(boxes):
        dbox(s, bx[i], ya, bw[i], hb, txt, fill=fills[i], edge=edges[i], size=13)
    for axx in arx:
        harrow(s, axx, ya + hb / 2 - 0.09, 0.45)
    bullets(s, [b0("The **encoder** maps the field to a latent: a pure state map ω(t) → z"),
                b0("The **predictor** advances the latent z(t) → z(t+1) from the latent history; it **is** the latent dynamics model you plan or learn against"),
                b0("A **visualisation decoder** is trained on the **frozen** encoder, never in the loss (figures only)"),
                b0("Compared at matched latent dimension against a reconstructive AE (Fukami) and POD")],
            0.7, 4.75, 12.0, 2.1, base=15)
    left_footer(s)
    notes(s, "Encoder is a pure state map; the predictor advances the latent. For "
             "control, an actuation channel is the natural next input.")
    return s


def slide_decoder():
    s = slide()
    header(s, "The visualisation decoder (figures only, never in the loss)", "Method")
    # small scheme: latent z -> LapFiLM decoder -> vorticity field
    ya = 1.95
    hb = 0.95
    bx = [1.45, 4.4, 9.2]
    bw = [2.4, 4.4, 2.7]
    arx = [3.95, 8.85]
    boxes = ["latent z\nd = 64",
             "LapFiLM decoder\n5-level Laplacian pyramid + FiLM\npixelshuffle upsampling · ~0.9M params",
             "vorticity field\nω (192 × 96)"]
    fills = [JP_FILL, BOX_FILL, BOX_FILL]
    edges = [JP_GRN, NAVY, NAVY]
    sizes = [13, 12, 13]
    for i, txt in enumerate(boxes):
        dbox(s, bx[i], ya, bw[i], hb, txt, fill=fills[i], edge=edges[i], size=sizes[i])
    for axx in arx:
        harrow(s, axx, ya + hb / 2 - 0.09, 0.45)
    bullets(s, [b0("Trained on the **frozen** encoder, **never in the JEPA loss**: it only renders a latent as a field so we can look at it"),
                b0("**Why not a plain CNN or ViT decoder?** A transposed-conv CNN under a pixel objective **blurs the wake** and leaves checkerboard artefacts; a ViT decoder is heavier and patch-blocky. The **Laplacian pyramid** builds the field coarse-to-fine (large-scale wake first), **FiLM** lets the 64-D latent steer every scale at just ~0.9M params, and **pixelshuffle** upsamples without checkerboards"),
                b0("**Loss** = region + Laplacian-pyramid + gradient + **spectral-amplitude** (small enstrophy / circulation terms): it matches structure **across scales and in the Fourier spectrum**, not pixelwise"),
                b0("**Why SSIM, not MSE?** MSE is pointwise and **rewards blur** (a smoothed field scores well but loses the vortices); **SSIM** compares local contrast and structure, crediting a correct wake *pattern*. Held-out **SSIM ≈ 0.73**")],
            0.7, 3.35, 12.0, 3.25, base=14)
    left_footer(s)
    notes(s, "The visualisation decoder is a SEPARATE stage on the frozen encoder, never in "
             "the JEPA objective. LapFiLM = 5-level Laplacian pyramid + FiLM modulation + "
             "pixelshuffle upsampling, ~0.9M params. Why this and not a plain CNN/ViT decoder: a "
             "transposed-conv CNN under pixel MSE blurs the wake and adds checkerboard artefacts, "
             "and a ViT decoder is heavier and patch-blocky; the pyramid synthesises coarse-to-fine "
             "(large-scale wake first), FiLM lets the 64-D latent modulate every scale cheaply, and "
             "pixelshuffle (sub-pixel conv) upsamples without checkerboards. Loss is a multi-term "
             "region + Laplacian-pyramid + gradient + spectral-amplitude objective (plus small "
             "enstrophy/circulation terms). Why SSIM not MSE: MSE is pointwise and rewards a blurred "
             "mean field, whereas SSIM scores local luminance/contrast/structure, so it tracks "
             "whether the wake PATTERN is right. Held-out SSIM ~0.73. Used only to visualise latents "
             "as fields in the rollout / recovery figures.")
    return s


def slide_model_detail():
    return s_fig_right(
        "Model detail: encoder and predictor", "Method",
        [b0("**Encoder** = a CNN stem (3 down-sampling stages) + a 6-layer Vision Transformer; the field becomes 288 patch tokens, and a learned [CLS] token is mapped by a small MLP (with BatchNorm) to the latent z (d = 32 / 64)"),
         b0("**Predictor** = a 6-layer autoregressive transformer (hidden 384, 16 heads, dropout 0.1): it reads z(t) and emits the next latent z(t+1). **Unconditioned: no gust parameters enter the predictor**"),
         b0("**Causal mask:** each step sees only past steps, never the future; **RoPE** encodes the frame ordering so the transformer knows what is earlier or later"),
         b0("**Training:** match the next latent (teacher forcing) + **scheduled sampling** (feed the model its own predictions over an 8-step rollout, so it tolerates its own errors) + **SIGReg** anti-collapse (keeps z informative without a decoder); AdamW, bf16, 20k steps"),
         b0("**Architecture inspired by the JEPA world-model line:** I-JEPA / V-JEPA, and the from-pixels LeWM (Maes et al. 2026) and LeJEPA (Balestriero and LeCun 2025)")],
        "predictor_detail.png", fl=7.3, fw=5.6, base=12.5,
        note="Moved from backup into the architecture block (after the decoder slide). Encoder = "
             "CNN stem + 6-layer ViT, [CLS]->MLP(BatchNorm)->z. Predictor = 6-layer AR transformer, "
             "unconditioned (no gust parameters). Causal mask + RoPE. Trained with teacher forcing + "
             "scheduled sampling (8-step) + SIGReg. Lineage: I-JEPA/V-JEPA, LeWM, LeJEPA.")


# ---- build ------------------------------------------------------------------
s_title()

# I. Motivation -> objective -> literature -> research question
_mot = slide()
header(_mot, "The control problem: gusts are fast, strong, nonlinear", "Motivation")
bullets(_mot,
        [b0("Small and micro air vehicles increasingly fly where the **gust speed is comparable to or above the flight speed** (urban canyons, building and ship wakes): gust ratio **G > 1** brings large, fast load transients classical models miss"),
         b0("A vortex strikes the wing: a **leading-edge vortex** forms and sheds, and the lift swings strongly (right: DNS, a moderate **G = -2** gust, C_L from about -1 to +3)"),
         b0("Active control, model-based or RL, needs a **fast, faithful forward model** it can trust under rollout"),
         b0("**This talk: what reduced state is worth planning against?**")],
        0.7, 1.95, 5.0, 5.0, base=15)
_mot.shapes.add_movie(str(FIG / "gust_motivation_anim.mp4"), Inches(5.95), Inches(1.75),
                      Inches(6.7), Inches(4.76),
                      poster_frame_image=str(FIG / "gust_motivation_poster.png"), mime_type="video/mp4")
left_footer(_mot)
notes(_mot, "Pure DNS, no model yet: motivate the difficulty. A moderate in-envelope gust (G = -2, "
            "D = 0.5, train set) on Fukami's canonical core; the lift swings by about 4 in C_L, in his "
            "recognisable range (baseline C_L 0.73 matches Fukami's 0.73). Pitch to a control audience: the "
            "bottleneck for model-based control / RL is a reduced dynamics model trustworthy under rollout.")

s_bullets(
    "Objective: a reduced state that keeps the wake", "Objective",
    [b0("The load transient is built by the **leading-edge vortex (LEV)** and the wake it leaves: the discriminating information is in the **wake**, not the integrated forces C_L, C_D"),
     b0("We quantify it by the **large-scale wake enstrophy**: a Gaussian scale split (Motoori & Goto 2019) at σ/c = 0.05 isolates the load-bearing large-scale vorticity (LEV + shear layer) from fine turbulence, then integrates ω_L² over the wake"),
     b0("**Objective:** find a reduced state that faithfully **encodes the wake** (the part reconstruction loses), not just the forces, and that is **recoverable from sparse sensors**"),
     b0("And **compare candidate states**, predictive vs reconstructive vs linear, under one matched protocol")],
    note="Objective slide (no results yet): the wake challenge, the large-scale wake enstrophy "
         "definition (Gaussian split), and the goal: a reduced state that ENCODES the wake (which "
         "reconstruction loses) and is recoverable from sensors. The predictive-vs-reconstructive "
         "Gaussian-split result comes later, in Results.")

s_question(
    "The opening for JEPA", "Literature and idea",
    [b0("Most flow ROMs, POD/DMD, autoencoders, are trained to **reconstruct the field**"),
     b1("reconstruction fixes the latent only up to a diffeomorphism, so its geometry is not constrained to be predictable"),
     b0("Alternative: learn a state that is **predictive by construction**, and plan against it")],
    "Which reduced state faithfully encodes the wake that reconstruction loses,\n"
    "and is recoverable from sparse sensors?",
    note="Set up the central question. The state should ENCODE the wake (which reconstruction-only "
         "states lose) and be recoverable from sensors. This is exactly the property a planner or RL "
         "agent needs to consume.")

# II. What a JEPA is, and our instantiation
s_fig_right(
    "Data", "Data",
    [b0("**DNS** (SOD2D, no subgrid-scale model), NACA 0012, α = 14°, **Re = 5000**"),
     b0("Taylor-vortex gusts parametrised by strength **G**, core diameter **D**, wall-normal offset **Y** (right: the 84-case envelope, isometric)"),
     b0("Impact-centred encounters; five observables: C_L, C_D, **wake enstrophy**, ±circulation"),
     b0("Held-out **test_b** (interpolation) and **test_c** (|G| = 4 extrapolation)")],
    "paramspace3d.png",
    cap="Training envelope in (G, D, Y): train, test_b interpolation, test_c |G| = 4 extrapolation, baseline.",
    fl=7.3, fw=5.7, base=16,
    note="Keep the DNS claim bounded (solver-resolution metadata pending). Emphasise the six "
         "observables, wake enstrophy being the one that separates the methods; the 3D envelope shows "
         "test_c (|G|=4) sitting far outside the training cloud.")

slide_jepa_concept()

s_bullets(
    "What is a JEPA?  (2 / 2): the recipe and why for control", "Method, primer",
    [b0("**Training, in one line:** minimise the **latent** prediction error plus an **anti-collapse term** so the encoder cannot cheat by mapping everything to a constant (no decoder, no pixel loss)"),
     b0("So the encoder **keeps only what predicts the future** (the control-relevant dynamics), and the predictor **gets no pixel-level shortcuts**"),
     b0("**Cheap to roll out:** planning happens in the d = 32-64 latent, two to three orders of magnitude cheaper than evolving the full field"),
     b0("Lineage: I-JEPA / V-JEPA (image, video); from-pixels LeWM, LeJEPA, PLDM, so far only gridworld and toy visual tasks"),
     b0("**Why for control:** a learned world-model that is **observable from sensors**; to our knowledge the **first end-to-end JEPA on a parametric flow**")],
    note="Second JEPA slide: recipe (latent loss + anti-collapse, no decoder), cheap latent "
         "rollout, lineage, and the control framing.", base=17)

slide_our_arch()

# The visualisation decoder (figures only, never in the JEPA loss)
slide_decoder()

# Encoder / predictor detail (moved here from backup to group the architecture)
slide_model_detail()

# Illustrated alternatives (Fukami style); compare with the box diagrams on slides 6 & 8
s_fig_below(
    "Illustrated alternative: the autoencoder route", "Method, alternative",
    [b0("Reconstructive route (Fukami-Taira lineage): a **CNN encoder**, an **MLP latent ξ** (which also feeds a lift head to C_L), and a **CNN decoder**; the loss compares the reconstructed field to the input in **pixel space**")],
    "fukami_ae.png",
    cap="An autoencoder reconstructs the field; the latent is shaped by pixel reconstruction.",
    ft=2.7, fh=3.6, base=16,
    note="Illustrated alternative to the box diagram (slide 6); keep whichever style you prefer.")

s_fig_below(
    "Illustrated alternative: the JEPA route", "Method, alternative",
    [b0("Same encoder, but a **predictor** advances the latent and the loss lives in **latent space**, with **no decoder** (an anti-collapse term keeps z informative)")],
    "fukami_jepa.png",
    cap="JEPA predicts the next latent in latent space, with no field reconstruction.",
    ft=2.7, fh=3.7, base=16,
    note="Illustrated alternative to the box architecture (slide 8); keep whichever style you prefer.")

# III. Protocol
s_fig_below(
    "Closure protocol", "Evaluation",
    [b0("**Representational closure (the headline):** probe each observable from the **encoded** latent at impact + 16, what the state *carries*, no rollout"),
     b0("**Forward (rollout) mode:** roll the predictor recursively to **H = 16** and probe the predicted latent; here we use it **qualitatively** (does the rolled latent track the wake?)"),
     b0("**Matched protocol:** same predictor architecture and same probe family, **trained / fitted separately per latent**, so differences are the encoder's"),
     b0("**Parameter-only floor:** can the **three gust numbers (G, D, Y) alone** predict each observable (a kernel-ridge fit)? Beating the floor proves the latent **carries flow state beyond the parameters**")],
    "eval_protocol.png",
    cap="Same probe family maps every encoder (JEPA / AE / POD) to the observables. Representational mode (encoded latent) is the headline; forward rollout is shown qualitatively.",
    ft=4.35, fh=2.1, base=14,
    note="Fairness is the whole point: only the encoder changes. Headline is REPRESENTATIONAL closure "
         "(probe the encoded latent); the forward rollout is shown qualitatively only. The parameter-only "
         "floor rules out 'the parameters did it'. Reported with bootstrap CIs, 3 encoder seeds, "
         "5-fold probe CV; held-out test_b / test_c.")

# IV. Results
s_fig_right(
    "Main result: the latent keeps the wake", "Result",
    [b0("**Representational closure:** probe the **encoded** latent (no rollout) for each observable at impact + 16; the question is what the unconditioned state *carries*"),
     b0("**Wake enstrophy is the discriminator: tf-no-c R² = 0.71** (the conditioned model 0.75; reconstructive and POD far lower, 0.06 and below)"),
     b0("So the **unconditioned** latent (no gust parameters anywhere) **encodes the wake that reconstruction smooths away**"),
     b0("Forces and circulations are read off cleanly too (C_L 0.88, Γ± 0.85 / 0.78, C_D 0.69)"),
     b0("**Wake enstrophy** E_w = ∫ ω_z² over the wake (x/c ∈ [0.5, 4], |y/c| ≤ 1): the **intensity of the wake's rotational structures** (LEV + shed vortices) that set the future load, the part reconstruction-only states lose")],
    "repr_closure.png",
    cap="Held-out representational R² at impact + 16 from the encoded latent (test_b). Left: per observable (tf-no-c). Right: wake-enstrophy R² across families.",
    note="Headline, reframed to REPRESENTATIONAL closure of the unconditioned tf-no-c latent. "
         "Probe the ENCODED latent (z_dns), no rollout: wake-enstrophy R^2 = 0.71 (conditioned 0.75), "
         "vs Fukami 0.06 / POD negative. I_y is excluded (the latent does not encode it; not "
         "comparable to the others). No 'rolled-out "
         "0.84/0.93' claim; the rollout is shown only qualitatively on the next slides.",
    fl=7.6, fw=5.4, base=16)

_anim = slide()
header(_anim, "The rolled latent qualitatively tracks the wake", "Result, qualitative")
bullets(_anim,
        [b0("The predictor **rolls the latent forward from impact**; each scalar is read off the rolled latent by a fixed linear probe"),
         b0("**Wake enstrophy** and **lift** qualitatively track the DNS truth (black) through impact and the lift dip"),
         b0("**Green = reconstruction** = the model's **own encoded latent decoded with NO rollout** (encode-then-decode of the true frame, the representational **ceiling**); it is **not** a separate reconstructive autoencoder"),
         b0("**Orange = prediction** = the **rolled latent** (the predictor advanced forward, then decoded)"),
         b0("A representative **low-error** encounter (G = +1.5), not the hardest case; shown qualitatively, no closure R² claimed here")],
        0.7, 2.05, 4.7, 4.6, base=15)
_anim.shapes.add_movie(str(FIG / "forecast_anim.mp4"), Inches(5.7), Inches(1.75),
                       Inches(6.9), Inches(4.98),
                       poster_frame_image=str(FIG / "forecast_poster.png"), mime_type="video/mp4")
left_footer(_anim)
notes(_anim, "Embedded MP4 (plays in PowerPoint; a still shows otherwise). Two observables "
             "(wake enstrophy and lift). GREEN reconstruction = the model's OWN encoded latent "
             "decoded with NO rollout (encode-then-decode of the true frame), i.e. the "
             "representational ceiling, NOT a separate reconstructive autoencoder. ORANGE "
             "prediction = the rolled latent (predictor advanced forward, then decoded). "
             "Representative low-error gust encounter (G = +1.5, the unconditioned tf-no-c "
             "rollout). Qualitative: the rolled latent tracks the wake; no closure R^2 is claimed.")

_field = slide()
header(_field, "The same rollout in physical space", "Result, physical space")
_ftf = textbox(_field, 0.7, 1.5, 11.9, 0.5)
para(_ftf, "Decode the rolled latent every frame: the rolled latent qualitatively tracks the wake as full vorticity fields, not just scalars.",
     size=17, color=INK, first=True)
_field.shapes.add_movie(str(FIG / "field_anim.mp4"), Inches(0.85), Inches(2.25),
                        Inches(11.6), Inches(3.06),
                        poster_frame_image=str(FIG / "field_poster.png"), mime_type="video/mp4")
caption(_field, "DNS truth vs reconstruction (JEPA encode→decode of the latent) vs prediction (rolled latent), "
                "G = +1.5. At impact (H = 0) prediction = reconstruction; they diverge with horizon. "
                "SSIM(prediction) 0.75 at impact, 0.63 at H = 16. The latent keeps the LEV and shear layer; "
                "fine-scale wake is not retained at d = 64.",
        0.85, 5.5, 11.6)
left_footer(_field)
notes(_field, "Physical-space companion to the scalar rollout: same encounter (G = +1.5, unconditioned "
              "tf-no-c), decoded with the frozen LapFiLM visualisation decoder (never part of the JEPA "
              "loss). Reconstruction is the representational ceiling; prediction is the rolled latent "
              "decoded. Qualitative: the rolled latent tracks the wake.")

s_fig_right(
    "It is the wake, not the forces", "Result",
    [b0("Each bar is **how well the future wake is forecast**: **dark = the full latent** (all 64 coordinates together), **light = the single best coordinate**"),
     b0("**Predictive JEPA:** the full latent (**0.84**) far exceeds the best single coordinate (**0.48**) → the wake is a **distributed code spread across many coordinates**"),
     b0("**Reconstructive / POD:** the best single coordinate is **already as good as the whole latent** → there is **no distributed / collective wake code** to find"),
     b0("**Forces** (C_L, C_D), by contrast, are carried **redundantly by every family**, in a single coordinate"),
     b0("And it clears the **parameter-only floor** → the latent, not the parameters, does the work")],
    "wake_code.png",
    cap="Future-wake skill: full latent (dark, all 64 coordinates) vs the single best coordinate (light). Predictive = distributed code; reconstructive / POD = one coordinate suffices.",
    note="Mechanistic reading of the main result. Each bar = how well the FUTURE WAKE is forecast; "
         "dark = the full latent (all 64 coordinates together), light = the single best coordinate. "
         "For the predictive JEPA the full latent (0.84) far exceeds the best single coordinate "
         "(0.48): the wake is a DISTRIBUTED code. For reconstructive / POD the best single "
         "coordinate is already as good as the whole latent, so there is NO distributed/collective "
         "wake code to find. Forces are carried redundantly by every family.")

s_fig_right(
    "The wake in physical space: the Gaussian scale split", "Result",
    [b0("**Gaussian scale split** (Motoori & Goto 2019): low-pass the vorticity at σ/c = 0.05 into a **large-scale** part (LEV + shear layer, carries the lift) and a fine part"),
     b0("Top: large-scale vorticity at impact+16 for the strongest test gust (G = -3, D = 1.5, Y = -0.1), simulation vs the **predictive** and **reconstructive** encode-then-decode reconstructions: predictive keeps the LEV/wake, reconstructive smooths it"),
     b0("Bottom: the **large-scale wake enstrophy** through the encounter; bands = mean ± 1 s.d. across test_b encounters"),
     b0("So 'it is the wake' is visible in **physical space**, not only in the scalar closure")],
    "scale_decomp.png",
    cap="Large-scale (σ/c = 0.05) wake vorticity at impact+16 (top) and large-scale wake enstrophy vs frame (bottom), by family.",
    note="Physical-space version of 'it is the wake': the predictive reconstruction keeps the large-scale "
         "wake, the reconstructive one smooths it. These are encode-then-decode reconstructions; bands are "
         "mean +/- 1 s.d. over test_b. Moved here from the opening so results follow the method.",
    base=14)

s_fig_right(
    "Latent coordinates group by physical function", "Result",
    [b0("Profile each of the 64 unconditioned-latent coordinates by its |Spearman| correlation with nine descriptors (gust G, forces, wake enstrophy, circulations Γ±, wake thickness, centroid)"),
     b0("Clustering the profiles recovers **functional groups**: **59 wake-vorticity** coordinates (three wake-dominated clusters) and **5 gust-forcing** coordinates (one force cluster)"),
     b0("Alone, the wake groups carry the wake at **≈ 0.5-0.8**, the forcing group at **0.16** (all 64 together: 0.84), the collective wake code seen earlier"),
     b0("Read **descriptively**, not as a causal claim")],
    "coord_groups.png",
    cap="Per-coordinate |ρ| with each descriptor; rows grouped into functional clusters (G1-G4), labelled with coord count and held-out wake skill.",
    note="Backs the 'it is the wake' story: the unconditioned latent organises coordinates by "
         "physical function (three wake-dominated clusters + one force cluster; full latent 0.84). "
         "Descriptive correlation only; the causal (SURD) part of this analysis is not in the paper.",
    fl=6.6, fw=6.4, base=15)

s_fig_right(
    "Controls: objective and supervision, not architecture", "Result, controls (conditioned-model control)",
    [b0("**(Conditioned-model control; no unconditioned variant run.)** **2 × 2:** objective {predictive, reconstructive} × architecture {CNN, CNN+ViT}, auxiliary heads matched"),
     b0("**Fair comparison:** every cell is rolled with the **same temporal predictor** (the matched closure predictor: same architecture and protocol, trained separately per family), so the **only thing that varies is the encoder's objective / architecture**"),
     b0("The **predictive objective wins at both architectures** (wake R² 0.46 / 0.45 vs 0.16 / 0.29) → **it is the objective, not the ViT**"),
     b0("And the **wake-observable head is a necessary ingredient**: removing it removes the gain"),
     b0("→ the gain needs the **predictive objective and wake supervision together**")],
    "controls_fairness.png",
    cap="Objective × architecture controls, three seeds per cell; every cell rolled with the same matched closure predictor, so only the encoder's objective / architecture varies.",
    note="Pre-empt the skeptic: isolate the cause before the mechanism. Fair-comparison setup: every "
         "cell is rolled with the SAME temporal predictor (the matched closure predictor, same "
         "architecture and protocol, trained separately per family), so the only thing that varies is "
         "the encoder's objective/architecture. The predictive objective wins at both architectures "
         "(it is the objective, not the ViT), and the wake-observable head is a necessary ingredient.")

s_fig_right(
    "Diagnostic 1: latent drift", "Mechanism",
    [b0("**Why this question:** a planner / RL agent queries the model at states reached by its **own rollout**, one-step error is not enough"),
     b0("**z_dns** = the latent **encoded directly from the simulation field** (ground truth); **z_markov** = the latent **produced by rolling the predictor forward**"),
     b0("The curve is their **relative distance vs horizon** (‖z_markov − z_dns‖ / ‖z_dns‖); **lower / flatter is better** (the rollout stays on the data manifold)"),
     b0("**Result:** the **unconditioned** rollout drift grows **gracefully** and **matches the conditioned model**; it stays on-manifold, exactly where the probes (and a planner) are valid")],
    "drift.png",
    cap="Rollout drift vs horizon (test_b): how far the rolled latent z_markov strays from the true encoded latent z_dns (relative distance). Lower / flatter = stays on-manifold.",
    note="Define the terms for the audience. z_dns = the latent encoded directly from the simulation "
         "field (ground truth); z_markov = the latent produced by rolling the predictor forward. The "
         "curve is their relative distance vs horizon, and LOWER/flatter is better (the rollout stays "
         "on the data manifold). The unconditioned rollout drift grows gracefully and matches the "
         "conditioned model. Figure now plots only tf-no-c, lstm-no-c, and the conditioned production "
         "JEPA (Fukami dropped: its smooth latent gives a misleadingly low drift).",
    base=14)

s_fig_right(
    "Diagnostic 2: topology of the encounter", "Mechanism",
    [b0("A gust encounter is a **cycle**: the flow loops and roughly returns to where it started"),
     b0("We ask whether the latent traces **one clean loop** or a **tangle of many**"),
     b0("**Persistent homology** is just a principled **loop-counter** on the latent path (it counts the loops that genuinely persist, not noise)"),
     b0("**Result:** the **unconditioned** JEPA latent makes **one loop** (median 1); the **reconstructive** latent **fragments into many** (median ~4); p ≈ 5×10⁻⁹"),
     b0("So the predictive state captures the encounter as a **single coherent cycle** (count of persistent H1 generators; Mann-Whitney test)")],
    "cycle.png",
    cap="The latent path of an encounter: the predictive latent traces one loop (median 1), the reconstructive latent many (median ~4). Persistent homology counts the loops.",
    note="Plain-language framing for a non-expert audience. A gust encounter is a cycle (the flow "
         "loops and roughly returns). We ask whether the latent traces ONE clean loop or a tangle of "
         "many. Persistent homology is a principled loop-counter on the latent path. Result: the "
         "unconditioned JEPA latent makes one loop (median 1); the reconstructive latent fragments "
         "into many (median ~4); p ~ 5e-9. So the predictive state captures the encounter as a single "
         "coherent cycle. (Technically: count of persistent H1 generators on z_dns over 42 test_b "
         "encounters; encounter-level Mann-Whitney p ~ 5e-9. The PCA panels are encode->decode, no "
         "predictor; the orbit need not visually close in the 120-frame window.)", base=14)

s_fig_below(
    "Decoded reconstructions", "Result, physical space",
    [b0("Decoded fields: predictive vs reconstructive vs POD vs DNS"),
     b0("The predictive decode is **blurrier pixel-wise** but preserves the **transported large-scale wake**; the reconstructive decode is sharp yet drifts under rollout")],
    "reconstructions.png", cap="Held-out reconstructions; the predictive objective trades pixel sharpness for transported structure.",
    ft=3.4, fh=3.2)

s_fig_below(
    "Sparse sensor placement", "Control relevance",
    [b0("Sparse wall-pressure sensors selected by **TCSI / qDEIM** vs uniform"),
     b0("**Model:** the pressure → latent map is a **KernelRidge (RBF)** regressor on the K-sensor × pre-impact-window vector; K = 2 / 4 / 8 / 16 (LSTM / KRR 5-fold-CV selected to guard small-sample overfitting)"),
     b0("**Metric — state-recovery R²:** the **fraction of the held-out latent's variance recovered from the K pressure taps** (1 = perfect, 0 = no better than the mean), for the **unconditioned tf-no-c latent** (no gust parameters anywhere)")],
    "sensor_placement.png", cap="Optimal sparse sensor placement on the airfoil surface; state-recovery R² is the fraction of the unconditioned latent's variance recovered from K taps.",
    ft=4.0, fh=2.9)

s_fig_right(
    "Flow recovered from sparse wall pressure", "Control relevance",
    [b0("**Model:** a **KernelRidge (RBF kernel)** regressor maps the K wall-pressure taps over a pre-impact window to the **impact-frame latent**; the frozen decoder then renders the field"),
     b0("**K = 8 taps** recover the leading-edge vortex and shear layer; **K = 2** coarsens but keeps the gross wake structure"),
     b0("Benchmarked against the **best-case decode** (decode of the **true simulation latent**, i.e. the ceiling decode) and the simulation"),
     b0("So the predictive state is **reconstructible from a few wall sensors**, a deployment-relevant observability result")],
    "flow_recovery.png",
    cap="Flow recovered from sparse wall pressure: simulation, best-case (ceiling) decode of the true simulation latent, and decode of the latent estimated from K = 8 and K = 2 taps (held-out).",
    note="Reframed to match the figure (flow from pressure). The 'best-case decode' column is the "
         "CEILING decode (decode of the true simulation-encoded latent), i.e. the best the decoder "
         "could do given a perfect latent. The predictive-vs-reconstructive wake tracking is already "
         "on the scale-decomposition slide.", base=16)

_sense = slide()
header(_sense, "Sparse-sensor state estimation in action", "Control relevance")
_stf = textbox(_sense, 0.7, 1.5, 11.9, 0.5)
para(_stf, "No predictor here: a per-frame **MLP** maps a causal 6-frame window of 8 wall-pressure taps to the latent, then the frozen decoder renders the field.",
     size=16, color=INK, first=True)
_sense.shapes.add_movie(str(FIG / "sensing_anim.mp4"), Inches(1.75), Inches(2.15),
                        Inches(9.8), Inches(3.59),
                        poster_frame_image=str(FIG / "sensing_poster.png"), mime_type="video/mp4")
caption(_sense, "DNS truth (green dots mark the 8 TCSI taps) vs the field recovered from wall pressure "
                "alone, on a held-out encounter. Held-out-encounter latent R² = 0.74, SSIM 0.63. "
                "Wall pressure fixes the near-body LEV and shear layer; the far wake is not "
                "surface-observable.",
        0.9, 5.85, 11.5)
left_footer(_sense)
notes(_sense, "Dynamic companion to the static flow-recovery slide: a causal 6-frame pressure window at "
              "8 TCSI taps -> MLP -> latent -> frozen decoder, applied frame by frame (unconditioned "
              "tf-no-c latent). This is state ESTIMATION, not forecasting; the JEPA predictor is not used here.")

# V. Control relevance
s_fig_right(
    "Control relevance: observable ahead of impact", "Control relevance",
    [b0("From a causal window of sparse wall pressure, recover the **impact-frame latent** at a **lead time before impact** (held-out test_b)"),
     b0("**Impact state** is recovered at **R² ≈ 0.88** right at impact and stays **above 0.83 out to 8 frames ahead** (kernel ridge; LSTM comparable)"),
     b0("**Impact lift** C_L is estimated with MAE rising gracefully from **≈ 0.38 at impact to ≈ 0.60 at 8 frames ahead** (LSTM lower, 0.29 to 0.49)"),
     b0("So the **unconditioned** state is **observable ahead of impact** from a few wall sensors: the ingredient a controller consumes")],
    "cl_inference_simple.png",
    cap="Recovered from sparse wall pressure at each lead time before impact (test_b): (a) impact-state R²; (b) impact-C_L MAE. Kernel ridge vs LSTM.",
    note="The venue hook, reframed to the unconditioned figure: from sparse wall pressure we recover the "
         "impact-frame latent ahead of impact (state R^2 ~0.88 at impact, >0.83 out to 8 frames) and "
         "estimate impact lift (MAE 0.38->0.60 over 0-8 frames). Observable ahead of impact, the "
         "ingredient a controller needs. No predictor-in-loop / oracle three-way claim here.",
    fl=7.5, fw=5.5, base=16)

s_fig_below(
    "Two forecast windows: early warning, then forecast", "Control relevance",
    [b0("We **roll the unconditioned JEPA's own predictor forward in latent space** and read each observable off the rolled latent with a **fixed linear probe**: we **predict the JEPA latent, then probe wake / lift from it**"),
     b0("**Wake:** a wide **±16-frame (≈ ±0.8 t/c)** window around impact (R² 0.7-0.9) -- the structural signature gives **early warning** and short-range forecast"),
     b0("**Lift C_L:** hard to anticipate (flat before impact) but **strongly forecastable after** (R² ≈ 0.8 to +12) -- the load is predictable **once impact hits**")],
    "forecast_windows.png",
    cap="Wake (left) and lift (right) R² vs frames relative to impact (test_b): x < 0 anticipation lead, x > 0 forecast horizon; tf-no-c and lstm-no-c.",
    note="Two forecast windows for the unconditioned predictor (markov rollout, wake/C_L linear probes on "
         "the rolled latent): wake is predictable in a roughly symmetric +/-16-frame window around impact "
         "(early warning); lift is asymmetric -- hard to anticipate (flat pre-impact) but strongly "
         "forecastable post-impact.",
    ft=3.7, fh=2.8, base=14)

# VI. Conclusions
s_bullets(
    "Conclusions", "Conclusions",
    [b0("A **fully unconditioned predictive latent** (no gust parameters anywhere) **representationally keeps the wake** (wake R² ≈ 0.71) that **force-only and reconstruction-only** states lose"),
     b0("That latent **traces a single clean cycle**, is **transport-consistent** within an encounter, and is **observable from sparse wall pressure**"),
     b1("**compact** (d = 32 ≈ d = 64, participation ratio ≈ 1.7)"),
     b0("The **rollout** is shown **qualitatively**: it stays on-manifold and tracks the wake, but we make **no forward-closure R² claim**"),
     b0("It is a **substrate** for model-based control and RL world-models, **not yet a validated controller**"),
     b0("**Next:** add an **actuation channel** (same machinery), **close the loop**, push to **3D observability** at |G| = 4, and carry the recipe to **other parametric flows**")],
    note="Single closing slide, reframed to the UNCONDITIONED latent. What we built: an unconditioned "
         "predictive latent that representationally keeps the wake (R^2 ~0.71), traces a single clean "
         "cycle, is transport-consistent and observable from sparse pressure. The rollout is qualitative; "
         "no forward-closed wake R^2 0.84 claim. Honest boundary: a substrate, not a validated controller. "
         "Outlook: actuation channel, close loop, 3D observability, other flows.")

# Acknowledgements
_ack = slide()
header(_ack, "Acknowledgements", "Thanks")
_atf = textbox(_ack, 0.7, 1.5, 11.9, 0.5)
para(_atf, "This work is a collaboration with:", size=18, color=INK, bold=True, first=True)

# author photo tiles
_people = [("photo_solera", "Alberto Solera-Rico", "INTA · UC3M", 0.50, 0.45),
           ("photo_miro", "Arnau Miró", "UPC · BSC", 0.50, 0.45),
           ("photo_lehmkuhl", "Oriol Lehmkuhl", "BSC", 0.33, 0.45)]
_cols = [2.55, 6.65, 10.75]
_ps = 1.7
for (_key, _name, _aff, _fx, _fy), _cx in zip(_people, _cols):
    _sp = square_photo(_key, focus_x=_fx, focus_y=_fy)
    if _sp:
        place_image(_ack, _sp, _cx - _ps / 2, 2.2, _ps, _ps)
    _ntf = textbox(_ack, _cx - 1.85, 4.0, 3.7, 0.75)
    para(_ntf, _name, size=15, color=NAVY, bold=True, align=PP_ALIGN.CENTER, first=True)
    para(_ntf, _aff, size=12, color=MUTE, align=PP_ALIGN.CENTER)

# funding (text left, BBVA + Red Leonardo logos right)
rect(_ack, 0.7, 5.05, 11.9, 1.75, fill=RGBColor(0xF3, 0xF5, 0xF9))
rect(_ack, 0.7, 5.05, 0.10, 1.75, fill=ACCENT)
_ftf = textbox(_ack, 1.05, 5.22, 8.2, 1.5)
para(_ftf, "Work produced with the support of a 2024 Leonardo Grant for Researchers and Cultural "
           "Creators, BBVA Foundation, project PREVENT (grant LEO24-2-15988). The Foundation takes "
           "no responsibility for the opinions, statements and contents of this project, which are "
           "entirely the responsibility of its authors.", size=13, color=INK, first=True)
logo(_ack, "bbva", 9.55, 5.3, 2.7, 0.62)
logo(_ack, "redleonardo", 9.55, 6.05, 2.7, 0.62)
left_footer(_ack)
notes(_ack, "Acknowledgements: author photos (Solera-Rico, Miró, Lehmkuhl) and the BBVA Foundation / "
            "Red Leonardo funding logos. Funding: 2024 Leonardo Grant, project PREVENT (LEO24-2-15988).")

# ---- backups ----------------------------------------------------------------
s_divider("Backup slides", "Supporting detail and ablations")

s_fig_right(
    "Transport geometry (optimal transport)", "Backup, mechanism",
    [b0("**Why:** Euclidean / pixel distance ignores **where structures move**; control cares about advection of the LEV and shear layer"),
     b0("**OT field distance** (after Tran, Yeh & Taira, JFM 2026): split vorticity into ± parts, transport each with **unbalanced Sinkhorn** OT and sum, the least work to rearrange one field into another"),
     b0("We did **not** train the latent to match OT (that paper does); this is a **post-hoc test**: per encounter, Spearman-correlate the **latent** distance matrix with the **OT** matrix (n = 42)"),
     b0("**Within-encounter:** the unconditioned JEPA **0.61** vs reconstructive **0.45**, order-preservation along the trajectory, **not** an isometry"),
     b0("**Honest caveat:** *pooled across* encounters the alignment **does not hold**, so it is **trajectory-local**, not a global latent metric")],
    "ot.png",
    cap="Per-encounter latent-vs-OT distance alignment (Shepard). OT field distance after Tran, Yeh & Taira (JFM 2026). Within-encounter Spearman 0.61 vs 0.45; pooled it does not hold.",
    note="Moved to backup: real but the weakest mechanism diagnostic (post-hoc, trajectory-local, "
         "pooled reverses). OT method = signed-vorticity split + unbalanced Sinkhorn (Tran, Yeh & Taira "
         "JFM 2026, eq 5); our code uses the entropic transport cost, not the debiased divergence, which "
         "does not change orderings.",
    base=13)

s_fig_right(
    "Dataset and protocol", "Backup",
    [b0("Partition **v2** (locked, sha256-anchored): 84 cases"),
     b0("**226** train encounters · **42** test_b · **24** test_c (|G| = 4)"),
     b0("ω_z pipeline v1: spatial mask + p99.99 clip + 3σ scale (train_std 3.55), no mean shift (preserves vorticity antisymmetry)"),
     b0("Reporting: bootstrap n = 2000 CIs · 3 encoder seeds · 5-fold probe CV")],
    "paramspace.png", cap="Training envelope in (G, D, Y).")

s_bullets(
    "Matched-capacity: d = 32 vs d = 64", "Backup (conditioned-model control)",
    [b0("**(Conditioned-model control; no unconditioned variant run.)** Halving the latent leaves the representation and **every mechanism diagnostic** intact"),
     b1("representational wake R² 0.74 (d=32) vs 0.75 (d=64); drift ratio 0.86 vs 0.85; OT 0.61 vs 0.63"),
     b0("Cost is **in-distribution forecast sharpness**: H = 16 wake closure 0.45 → 0.21"),
     b0("Participation ratio ≈ 1.7, leading PC ≈ ¾ of variance → the effective dimension is a handful")])

s_bullets(
    "Seed variance", "Backup (conditioned-model control)",
    [b0("**(Conditioned-model control; no unconditioned variant run.)** Three encoder seeds per cell; wake-closure means reported ± standard deviation"),
     b0("The reconstructive CNN+ViT cell has **large seed variance (±0.27)**, consistent with the drift mechanism: without a forward-predictable geometry, closure is unstable across initialisations"),
     b0("Predictive cells are tight (±0.03–0.06)")])

s_fig_below(
    "Parameter observability from the latent", "Backup",
    [b0("z → (G, D, Y) probe R² from the rolled-out latent, K = 8 pre-impact sensors"),
     b0("test_b: **G 0.46 · D 0.80 · Y 0.10**, gust strength and core diameter recoverable; cross-stream offset Y is marginal"),
     b0("test_c (|G| = 4): negative, the 3D observability boundary of a single mid-plane slice")],
    "pressure_observability.png", cap="Gust parameters implicit in the latent on held-out cases.",
    ft=3.7, fh=2.6)

s_bullets(
    "SSIM convention", "Backup",
    [b0("Reconstruction SSIM uses the **Wang convention** (K1 = 0.01, K2 = 0.03) on pipeline-normalised ω"),
     b0("Data range **L = 2 · global p99.9(|target|) ≈ 8.31** (split v2)"),
     b0("Decoder test_a SSIM ≈ **0.73** at this convention")])

s_bullets(
    "Parameter-only floor", "Backup",
    [b0("KRR-RBF on (G, D, Y) → observable at impact; train / test_b / test_c R²"),
     b1("on train, (G, D, Y) interpolates C_L and C_D well (226 points in 3-D)"),
     b0("On **test_b** the floor **collapses** across observables; on **test_c** it is negative for 5 of 6"),
     b0("→ JEPA's generalisation is **not** explained by the gust parameters alone")])

s_fig_right(
    "Phase–amplitude reading", "Backup",
    [b0("Phase–amplitude decomposition of the encounter cycle in the predictive latent"),
     b0("The sharp **2π → 0 step** in panel (b) is a **phase wrap** (θ is a cyclic angle on the orbit; it resets at the branch cut, which falls near impact), **not a physical jump**"),
     b0("Connects the predictor to the **sensitivity-function control** designed for these flows"),
     b0("An actuation channel is the natural next input for closed-loop control")],
    "phase_amplitude.png", fl=7.9, fw=5.0,
    cap="The 2π → 0 step in panel (b) is a phase wrap (cyclic angle resetting at the branch cut near impact), not a physical discontinuity.",
    note="Phase-amplitude reading of the encounter cycle. Note for the audience: the abrupt "
         "2*pi -> 0 step in panel (b) is a phase WRAP (theta is a cyclic angle on the orbit, it "
         "resets at the branch cut, which falls near impact), not a physical jump.")

# ---- page numbers (skip the title) -----------------------------------------
for i, sl in enumerate(prs.slides):
    if i == 0:
        continue
    tf = textbox(sl, 12.35, 7.06, 0.75, 0.3)
    para(tf, str(i + 1), size=9.5, color=MUTE, align=PP_ALIGN.RIGHT, first=True)

prs.save(str(OUT))
print(f"wrote {OUT}  ({len(prs.slides)} slides)")
